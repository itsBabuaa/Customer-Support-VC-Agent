"""Hybrid RAG: FAISS (semantic) + BM25 (keyword) with RRF re-ranking."""

import re
import os
import math
import logging
import numpy as np
from pathlib import Path
from collections import Counter

import faiss
from openai import OpenAI

logger = logging.getLogger("voice-agent.rag")

# ── File readers ───────────────────────────────────────────


def _read_txt(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "shift_jis", "cp932"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, Exception):
            continue
    return ""


def _read_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        logger.warning("Failed to read docx %s: %s", path.name, e)
        return ""


def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
        return "\n".join(texts)
    except Exception as e:
        logger.warning("Failed to read pptx %s: %s", path.name, e)
        return ""


def _read_pdf(path: Path) -> str:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        logger.warning("Failed to read pdf %s: %s", path.name, e)
        return ""


def _read_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None]
                if vals:
                    texts.append(" | ".join(vals))
        return "\n".join(texts)
    except Exception as e:
        logger.warning("Failed to read xlsx %s: %s", path.name, e)
        return ""


READERS = {
    ".txt": _read_txt, ".md": _read_txt,
    ".docx": _read_docx, ".pptx": _read_pptx,
    ".pdf": _read_pdf, ".xlsx": _read_xlsx,
}


# ── Chunker ────────────────────────────────────────────────

def _smart_chunk(text: str) -> list[str]:
    """Extract Q&A pairs first, then section/paragraph chunks."""
    chunks: list[str] = []

    # Pass 1: Q&A pairs
    qa_re = re.compile(
        r'(?:\d+\.\s*)?(?:\*{0,2})Q\d*[:：]\s*(.+?)(?:\*{0,2})\s*'
        r'(?:\*{0,2})A\d*[:：]\s*(.+?)(?=(?:\d+\.\s*)?(?:\*{0,2})Q\d*[:：]|\Z)',
        re.DOTALL | re.IGNORECASE,
    )
    for qm, am in qa_re.findall(text):
        q = re.sub(r'\*{1,2}', '', qm).strip()
        a = re.sub(r'\*{1,2}', '', am).strip()
        q = re.sub(r'\*?EN:.*$', '', q, flags=re.MULTILINE).strip()
        a_lines = []
        for line in a.split('\n'):
            line = line.strip()
            if line.startswith('*EN:') or line.startswith('EN:'):
                a_lines.append(re.sub(r'^\*?EN:\s*', '', line))
            elif not line.startswith('*') or len(line) > 5:
                a_lines.append(line)
        a_clean = ' '.join(a_lines).strip()
        if len(q) > 5 and len(a_clean) > 5:
            chunks.append(f"Q: {q}\nA: {a_clean}"[:600])

    # Pass 2: Section/paragraph chunks
    sections = re.split(r'\n#{1,3}\s+', text)
    for section in sections:
        section = section.strip()
        if len(section) < 40:
            continue
        paras = re.split(r'\n\s*\n', section)
        buf = ""
        for p in paras:
            p = p.strip()
            if not p or len(p) < 15:
                continue
            if len(buf) + len(p) < 500:
                buf += "\n" + p
            else:
                if len(buf.strip()) > 30:
                    chunks.append(buf.strip()[:600])
                buf = p
        if len(buf.strip()) > 30:
            chunks.append(buf.strip()[:600])

    # Deduplicate
    seen = set()
    unique = []
    for c in chunks:
        key = c[:80].lower()
        if key not in seen:
            seen.add(key)
            unique.append(c)
    return unique if unique else [text[:600]]


# ── Main RAG class ─────────────────────────────────────────

class KnowledgeRAG:
    """Hybrid retriever: FAISS (OpenAI embeddings) + BM25 with RRF."""

    EMBED_MODEL = "text-embedding-3-small"
    EMBED_DIM = 1536

    _SPLIT_RE = re.compile(
        r'[\s\-_/\\()（）【】「」『』、。，．・：；!?！？\*#\[\]{}|<>~`\'\"]+',
    )
    _STOP = frozenset({
        "the", "is", "are", "was", "were", "a", "an", "of", "in", "to",
        "for", "and", "or", "on", "at", "by", "it", "be", "as", "do",
        "if", "no", "not", "this", "that", "with", "from", "what", "how",
    })

    def __init__(self, *sources: str | Path):
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.chunks: list[str] = []
        self.faiss_index: faiss.IndexFlatIP | None = None
        self.idf: dict[str, float] = {}
        self.tf_vectors: list[Counter] = []
        self.avg_dl: float = 0.0
        self.chunk_lengths: list[int] = []

        raw_text = ""
        for src in sources:
            src = Path(src)
            if src.is_file():
                raw_text += "\n" + self._read(src)
            elif src.is_dir():
                for f in sorted(src.rglob("*")):
                    if f.is_file() and f.suffix.lower() in READERS:
                        raw_text += "\n" + self._read(f)

        self.chunks = _smart_chunk(raw_text)
        logger.info("Created %d chunks, building indexes...", len(self.chunks))
        self._build_faiss_index()
        self._build_bm25_index()
        logger.info("Hybrid retriever ready: %d chunks", len(self.chunks))

    def _read(self, path: Path) -> str:
        reader = READERS.get(path.suffix.lower())
        return reader(path) if reader else ""

    def _embed(self, texts: list[str]) -> np.ndarray:
        resp = self.client.embeddings.create(model=self.EMBED_MODEL, input=texts)
        arr = np.array([i.embedding for i in resp.data], dtype=np.float32)
        norms = np.linalg.norm(arr, axis=1, keepdims=True)
        norms[norms == 0] = 1
        return arr / norms

    def _build_faiss_index(self):
        if not self.chunks:
            return
        embeddings = self._embed(self.chunks)
        self.faiss_index = faiss.IndexFlatIP(embeddings.shape[1])
        self.faiss_index.add(embeddings)

    def _tokenize(self, text: str) -> list[str]:
        return [t for t in self._SPLIT_RE.split(text.lower())
                if len(t) > 1 and t not in self._STOP]

    def _build_bm25_index(self):
        n = len(self.chunks)
        if n == 0:
            return
        self.tf_vectors = []
        df: Counter = Counter()
        for chunk in self.chunks:
            tf = Counter(self._tokenize(chunk))
            self.tf_vectors.append(tf)
            for term in tf:
                df[term] += 1
        self.idf = {
            t: math.log((n - f + 0.5) / (f + 0.5) + 1.0)
            for t, f in df.items()
        }
        self.chunk_lengths = [sum(tf.values()) for tf in self.tf_vectors]
        self.avg_dl = sum(self.chunk_lengths) / n if n else 1.0

    def _bm25_score(self, tokens: list[str], idx: int) -> float:
        k1, b = 1.5, 0.75
        tf = self.tf_vectors[idx]
        dl = self.chunk_lengths[idx]
        score = 0.0
        for qt in tokens:
            freq = tf.get(qt, 0)
            if freq and qt in self.idf:
                score += self.idf[qt] * (freq * (k1 + 1)) / (
                    freq + k1 * (1 - b + b * dl / self.avg_dl))
        return score

    def _faiss_retrieve(self, query: str, top_k: int) -> list[tuple[int, float]]:
        if self.faiss_index is None:
            return []
        scores, indices = self.faiss_index.search(self._embed([query]), top_k)
        return [(int(i), float(s)) for i, s in zip(indices[0], scores[0]) if i >= 0]

    def _bm25_retrieve(self, query: str, top_k: int) -> list[tuple[int, float]]:
        tokens = self._tokenize(query)
        if not tokens:
            return []
        scored = [(i, self._bm25_score(tokens, i)) for i in range(len(self.chunks))]
        scored.sort(key=lambda x: x[1], reverse=True)
        return [(i, s) for i, s in scored[:top_k] if s > 0]

    @staticmethod
    def _rrf_merge(*ranked_lists: list[tuple[int, float]], k: int = 60) -> list[int]:
        fused: dict[int, float] = {}
        for ranked in ranked_lists:
            for rank, (idx, _) in enumerate(ranked):
                fused[idx] = fused.get(idx, 0.0) + 1.0 / (k + rank + 1)
        return [i for i, _ in sorted(fused.items(), key=lambda x: x[1], reverse=True)]

    def search(self, query: str, top_k: int = 3) -> str:
        if not self.chunks:
            return "No knowledge base loaded."
        faiss_r = self._faiss_retrieve(query, top_k * 2)
        bm25_r = self._bm25_retrieve(query, top_k * 2)
        merged = self._rrf_merge(faiss_r, bm25_r)
        results = [self.chunks[i] for i in merged[:top_k]]
        return "\n---\n".join(results) if results else "No relevant information found."
